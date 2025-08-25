"""
brochure_data = [
    {
        "title": "Welcome",
        "description": "We are pioneers in sustainable engineering.",
        "image": "images/cat_on_sofa.jpg",
        "features": {
            "Founded": "1999",
            "Headquarters": "Oslo, Norway",
            "Focus": "Green infrastructure"
        }
    },
    {
        "title": "Our Mission",
        "description": "To build a better world using eco-friendly technologies.",
        "image": "images/cat_on_sofa.jpg",
        "features": {
            "Core Value": "Sustainability",
            "Approach": "Collaborative",
            "Impact": "Global"
        }
    },
    {
        "title": "Services",
        "description": "We provide consulting, design, and execution.",
        "image": "images/cat_on_sofa.jpg",
        "features": {
            "Consulting": "Feasibility & Planning",
            "Design": "Modern & Compliant",
            "Execution": "Efficient & On-time"
        }
    },
    {
        "title": "Our Team",
        "description": "Led by innovators and experts in the field.",
        "image": "images/cat_on_sofa.jpg",
        "features": {
            "Employees": "200+",
            "Experts": "50+ Engineers",
            "Offices": "15 Worldwide"
        }
    },
    {
        "title": "Contact Us",
        "description": "Let’s collaborate on your next project!",
        "image": "images/cat_on_sofa.jpg",
        "features": {
            "Email": "contact@greenbuild.com",
            "Phone": "+47 123 456 78",
            "Website": "greenbuild.com"
        }
    }
]
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import json
import re





def clean_text(text):
    # Remove markdown bold **text**
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    
    # Remove underscores _text_
    text = re.sub(r"_([^_]+)_", r"\1", text)

    # Remove backticks `code`
    text = re.sub(r"`([^`]*)`", r"\1", text)

    # Remove markdown headers like ## Heading
    text = re.sub(r"^#{1,6}\s*", "", text)

    # Remove triple quotes
    text = text.replace('"""', '')

    # Remove single/double quotes (optional)
    text = text.replace('"', '')
    text = text.replace("'", '')

    # Remove hyphens at the beginning of lines or bullets
    text = re.sub(r"^\s*-\s*", '', text, flags=re.MULTILINE)

    # Final strip
    return text.strip()
# Function to add text and image to slide
def populate_slide(slide, data):
    # Title placeholder
    title_placeholder = slide.shapes.title
    if title_placeholder:
        title_placeholder.text = data["title"]
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)  # Adjust title size here
    # Content placeholder (usually idx 1, depends on layout)

    content_placeholder = None
    shapes=[]
    for shape in slide.placeholders:
        phf = shape.placeholder_format
        if phf.type == 2:  # BODY placeholder type
            content_placeholder = shape
            shapes.append(shape)
    content_placeholder = shapes[1]
    if content_placeholder:
        content_placeholder.text = data["description"]
        for paragraph in content_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)  # Adjust text size here


    # Add image to the slide
    left = Inches(6.75)
    top = Inches(0.35)
    height = Inches(6.5)
    slide.shapes.add_picture(data["image"], left, top, height=height)

def populate_slide_(slide, data):
    # Set title (assumes slide has a title placeholder)
    title_placeholder = slide.shapes.title
    if title_placeholder:
        title_placeholder.text = data["title"]

    # Add image (assumes enough space, adjust position as needed)
    left = Inches(4)
    top = Inches(1.5)
    height = Inches(3.5)
    slide.shapes.add_picture(data["image"], left, top, height=height)

    # Add description
    txBox = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(3), Inches(1))
    tf = txBox.text_frame
    tf.text = data["description"]


    # Add key-value features
    left = Inches(5.5)
    top = Inches(2.5)
    for key, value in data["features"].items():
        p = tf.add_paragraph()
        p.text = f"{key}: {value}"
        p.level = 0

def create_brochure(id):
    script_path = f"./brochures/{id}/brochure_object.json"
    output_file_path = f"./brochures/{id}/presentation.pptx"
    with open(script_path, "r", encoding="utf-8") as f:
        data = json.load(f)
        slides_raw = data.get("slides", {})
    brochure_data = []
    for i in sorted(slides_raw, key=lambda x: int(x)):
        slide = slides_raw[i]
        title = clean_text(slide.get("title", ""))
        description = clean_text(slide.get("description", ""))
        image = slide.get("image", "")
        brochure_data.append({
            "title": title,
            "description": description,
            "image": image
        })
    # Load template
    prs = Presentation("./assets/template.pptx")
    
    # Clear existing slides if needed (optional)
    # del prs.slides[1:]
    # Create 5 slides with content
    for i, data in enumerate(brochure_data):
        slide_layout = prs.slide_layouts[12]  # Choose a simple layout with title + content
        slide = prs.slides.add_slide(slide_layout)
        populate_slide(slide, data)

    # Save new brochure
    prs.save(output_file_path)
    print(f"✅ Brochure created: {output_file_path}")

if __name__ == "__main__":
    #id = 7321484248304160
    #id = 7917787685442408
    id = 1167470051692233
    create_brochure(id)
